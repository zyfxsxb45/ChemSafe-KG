"""
ChemSafe-KG 期末汇报PPT v10 — 深度视觉优化
设计语言: 极简白/卡片投影/克制配色/精准对齐
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ═══════════ 配色 ═══════════
BG      = RGBColor(0xFF,0xFF,0xFF)
ACCENT  = RGBColor(0x2B,0x5F,0xA6)   # 稳重蓝
RED     = RGBColor(0xC0,0x39,0x2B)    # 暗红
GREEN   = RGBColor(0x1E,0x84,0x45)    # 暗绿
ORANGE  = RGBColor(0xD3,0x54,0x00)    # 暗橙
PURPLE  = RGBColor(0x6C,0x34,0x80)    # 暗紫
DARK    = RGBColor(0x2C,0x3E,0x50)    # 深灰蓝
TEXT    = RGBColor(0x2C,0x3E,0x50)    # 主文字
GREY    = RGBColor(0x95,0xA5,0xA6)    # 灰色
LIGHT   = RGBColor(0xF8,0xF9,0xFA)    # 卡片底色
BORDER  = RGBColor(0xE0,0xE4,0xE8)    # 卡片边框
WHITE   = RGBColor(0xFF,0xFF,0xFF)
SHADOW  = RGBColor(0xD5,0xDB,0xDF)    # 投影色

CHART_DIR = r"D:\课程文件\大二下\数据库技术及应用\大作业\期末汇报材料\charts"
OUT_PATH  = r"D:\课程文件\大二下\数据库技术及应用\大作业\期末汇报材料\ChemSafe-KG_期末汇报.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG

def add_dots(slide, section, total=5):
    for i in range(total):
        x = Inches(0.4 + i * 0.35)
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(0.2), Inches(0.16), Inches(0.16))
        d.fill.solid()
        d.fill.fore_color.rgb = ACCENT if i == section - 1 else BORDER
        d.line.fill.background()

def add_title(slide, text, top=Inches(0.35)):
    """Title: 42pt bold + thin accent line below"""
    tb = slide.shapes.add_textbox(Inches(0.7), top, Inches(11.9), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = DARK
    # accent line
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), top + Inches(0.82), Inches(2.0), Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()

def add_shadow(slide, left, top, w, h):
    """Subtle shadow under a card"""
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.03), top + Inches(0.03), w, h)
    shadow.fill.solid(); shadow.fill.fore_color.rgb = SHADOW
    shadow.line.fill.background()
    # Make it look soft by lowering opacity via XML
    shadow.shadow.inherit = False

def add_card(slide, left, top, w, h, title, body, accent):
    """Card with shadow + top accent line + clean typography"""
    add_shadow(slide, left, top, w, h)
    # card bg
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BORDER; s.line.width = Pt(0.75)
    # top accent bar
    ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.05))
    ab.fill.solid(); ab.fill.fore_color.rgb = accent; ab.line.fill.background()
    # title
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), w - Inches(0.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = accent
    # body
    tb2 = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.65), w - Inches(0.5), h - Inches(0.85))
    tf = tb2.text_frame; tf.word_wrap = True
    lines = body.split('\n')
    n = len(lines)
    sz = Pt(15) if n <= 6 else Pt(14) if n <= 10 else Pt(13)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = sz; p.font.color.rgb = TEXT
        p.space_after = Pt(1)

def add_section_page(slide, num, title, subtitle=""):
    """Full-page divider: giant number + title + accent line"""
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(4))
    tf = tb.text_frame
    # number
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(130); p.font.bold = True; p.font.color.rgb = ACCENT
    # title
    p2 = tf.add_paragraph(); p2.text = title
    p2.font.size = Pt(42); p2.font.color.rgb = DARK; p2.space_before = Pt(8)
    # accent line
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.0), Inches(4.0), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    # subtitle
    if subtitle:
        p3 = tf.add_paragraph(); p3.text = subtitle
        p3.font.size = Pt(17); p3.font.color.rgb = GREY; p3.space_before = Pt(16)

def add_metrics(slide, top, items):
    """Metric cards row: big number + small label"""
    n = len(items); gap = Inches(0.15)
    w = (Inches(12.2) - gap * (n - 1)) / n
    for i, (num, label, color) in enumerate(items):
        left = Inches(0.55) + i * (w + gap)
        add_shadow(slide, left, top, w, Inches(0.95))
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.95))
        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = BORDER; bg.line.width = Pt(0.75)
        # top color dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + w/2 - Inches(0.08), top + Inches(0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        # number
        tb = slide.shapes.add_textbox(left, top + Inches(0.28), w, Inches(0.38))
        tb.text_frame.paragraphs[0].text = num
        tb.text_frame.paragraphs[0].font.size = Pt(26); tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = color; tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # label
        tb2 = slide.shapes.add_textbox(left, top + Inches(0.63), w, Inches(0.25))
        tb2.text_frame.paragraphs[0].text = label
        tb2.text_frame.paragraphs[0].font.size = Pt(10); tb2.text_frame.paragraphs[0].font.color.rgb = GREY
        tb2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_bottom(slide, text, top=Inches(6.65)):
    tb = slide.shapes.add_textbox(Inches(0.7), top, Inches(11.9), Inches(0.4))
    tb.text_frame.paragraphs[0].text = text
    tb.text_frame.paragraphs[0].font.size = Pt(14); tb.text_frame.paragraphs[0].font.color.rgb = GREY
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def ns(dots=True, sec=0):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    if dots: add_dots(s, sec)
    return s


# ══════════════════ P1 Cover ══════════════════
s = ns(dots=False)
tb = s.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.2)).text_frame
tb.paragraphs[0].text = "ChemSafe-KG"
tb.paragraphs[0].font.size = Pt(68); tb.paragraphs[0].font.bold = True; tb.paragraphs[0].font.color.rgb = ACCENT
tb.add_paragraph().text = "化工安全事故知识图谱的构建与效果验证"
tb.paragraphs[1].font.size = Pt(28); tb.paragraphs[1].font.color.rgb = DARK; tb.paragraphs[1].space_before = Pt(6)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.0), Inches(3.5), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
tb2 = s.shapes.add_textbox(Inches(1.2), Inches(4.4), Inches(10.9), Inches(2.5)).text_frame
tb2.paragraphs[0].text = "数据库技术及应用 · 期末汇报 · 2026年6月"
tb2.paragraphs[0].font.size = Pt(15); tb2.paragraphs[0].font.color.rgb = GREY
tb2.add_paragraph().text = "翟彝凡  余亮阳  赵乐毅  ·  清华大学化学工程系"
tb2.paragraphs[1].font.size = Pt(13); tb2.paragraphs[1].font.color.rgb = GREY; tb2.paragraphs[1].space_before = Pt(4)

# ══════════════════ P2 研究动机 ══════════════════
s = ns(sec=1)
add_title(s, "研究动机")
add_card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(3.5), "初始假设",
"知识图谱能提升LLM在化工安全问答中的表现\n\n出发点: LLM回答时会引用训练数据中的信息\n但这些信息无法追溯来源\n化工安全领域 不可追溯的回答存在风险\n\n我们的假设: 用事故报告建知识图谱\n让LLM只引用图谱里的内容\n会得到更可控的结果", ACCENT)
add_card(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(3.5), "验证方式",
"构建完整系统来检验这个假设\n\n1,300份事故报告 自动抽取因果链\n建成6,976节点 23,111边的知识图谱\n\n设计三组对照实验\n同一模型 唯一变量是检索方式\n比较回答的可追溯性和约束效果\n\n顺带做了多源数据融合和可视化分析", GREEN)
add_metrics(s, Inches(5.3),
[("1,579", "事故记录", RED), ("6,976", "KG节点", ACCENT),
 ("23,111", "关系边", GREEN), ("1947–2026", "时间跨度", ORANGE)])

# ══════════════════ P3 研究问题 ══════════════════
s = ns(sec=1)
add_title(s, "研究问题")
add_card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(2.5), "要回答什么",
"1 能否用LLM从非结构化的中文事故报告中\n  自动构建可用的知识图谱\n\n2 建成之后 这个图谱对LLM问答\n  到底有没有实际帮助\n\n3 如果有帮助 在什么条件下最有效\n  如果没帮助 原因是什么", ACCENT)
add_card(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.5), "挑战",
"中文化工报告长句多 术语密 格式不统一\n传统规则方法难以覆盖\n\nLLM抽取的准确性需要系统性评估\n不是看几个例子就能下结论\n\n三个数据源格式完全不同\n整合需要大量清洗对齐\n1,300份报告无法人工标注", RED)
add_bottom(s, "核心问题: 知识图谱+LLM对化工安全事故问答到底有没有用？", Inches(4.4))

# ══════════════════ P4 相关工作 ══════════════════
s = ns(sec=1)
add_title(s, "相关工作")
add_card(s, Inches(0.55), Inches(1.4), Inches(3.85), Inches(4.7), "传统统计分析",
"关文玲, 蒋军成 (2007)\n安全科学学报\n\n统计2001-2006化工火灾爆炸事故\n发现反应釜和储罐最高发\n\n方法: 统计频次 画分布图\n\n局限\n只能告诉你什么设备出事多\n说不清设备是怎么出的事\n做的是频次统计 不是因果追踪", GREY)
add_card(s, Inches(4.7), Inches(1.4), Inches(3.85), Inches(4.7), "KG工业化尝试",
"事理图谱方向 (2018年前后)\n关注事件间因果演化关系\n\n方法: 人工编写抽取规则\n例如文本含爆炸和温度升高就建关系\n\n局限\n化工事故类型多 规则写不全\n换数据源 规则要重调\n建图慢 贵 依赖人工", ORANGE)
add_card(s, Inches(8.85), Inches(1.4), Inches(3.85), Inches(4.7), "LLM + GraphRAG",
"微软 GraphRAG (2024)\n从非结构化文本自动提取实体关系\n\nGRAG-ProSafe QAS (2026)\n钢铁行业 198份英文报告\n1,637节点 2,285边\n\n未解决的问题\n英文场景 中文语境不同\n未和纯LLM做过严格对比\n未系统评估抽取准确率", ACCENT)

# ══════════════════ P5 分隔 ══════════════════
s = ns(dots=False)
add_section_page(s, "01", "系统构建", "架构 / 抽取 / 检索 / 存储 / 数据")

# ══════════════════ P6 架构 ══════════════════
s = ns(sec=2)
add_title(s, "五层系统架构")
img = os.path.join(CHART_DIR, "08_architecture.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.6))
else:
    add_card(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.5), "架构",
    "Web应用 · 检索 · 存储 · 抽取 · 数据采集", ACCENT)

# ══════════════════ P7 Prompt Chain ══════════════════
s = ns(sec=2)
add_title(s, "知识抽取: Prompt Chain")
add_card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(5.0), "8条迭代规则",
"事件原子化: 每个实体描述单一概念 不超过15字\n人员操作分离: 操作工误开阀门必须拆为\n  Equipment + Abnormal Condition\nEquipment/Material优先识别\n5种实体类型 × 3种关系类型\nFew-shot示例: 丙烯腈储罐爆炸完整因果链\nJSON三级容错: 空响应重试 代码块清理 正则强提\n200线程并发 DeepSeek v4-flash\n迭代十几版 成功率99%以上", ACCENT)
add_card(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.0), "Few-shot 示例",
"Equipment: 冷却水循环泵\n  involves Material: 丙烯腈\n  leads_to Abnormal: 冷却水中断\n    leads_to Abnormal: 储罐升温\n      leads_to Consequence: 聚合爆炸\n        mitigated_by Mitigation: 双路冷却\n\n示例同时展示正向格式和4组常见错误反例\nLLM对示例的模仿效果远好于对抽象规则的遵循", GREEN)

# ══════════════════ P8 三层匹配 ══════════════════
s = ns(sec=2)
add_title(s, "三层实体匹配")
add_card(s, Inches(0.55), Inches(1.4), Inches(3.85), Inches(3.8), "L1 精确匹配",
"jieba分词 直接比对KG节点名\n置信度 0.6 - 1.0\n\n权重 ×3\n最高优先级\n\n解决直接匹配问题\n用户说液氯 图谱有氯气", ACCENT)
add_card(s, Inches(4.7), Inches(1.4), Inches(3.85), Inches(3.8), "L2 关键词命中",
"jieba分词 与实体名交叉对比\n命中数计数\n\n权重 ×1\n辅助信号\n\n多个关键词同时命中\n说明实体相关性较高", ORANGE)
add_card(s, Inches(8.85), Inches(1.4), Inches(3.85), Inches(3.8), "L3 嵌入语义",
"sentence-transformers\n470MB多语言模型\n实体名清洗 + 自适应阈值\n\n权重 ×3\n解决同义词匹配\n液氯和氯气嵌入距离很近", GREEN)
add_bottom(s, "融合: score = L1×3 + L2×1 + L3×3 + 有出边奖励×2", Inches(5.6))

# ══════════════════ P9 Graph RAG ══════════════════
s = ns(sec=2)
add_title(s, "Graph RAG: 因果路径约束生成")
add_card(s, Inches(0.55), Inches(1.4), Inches(3.85), Inches(2.8), "按检索路径回答",
"基于因果路径生成答案\n不能添加推测性内容\n不能引用图谱中不存在的实体\n\n每条路径对应一个因果关系\n答案严格绑定在路径上", RED)
add_card(s, Inches(4.7), Inches(1.4), Inches(3.85), Inches(2.8), "路径不足时拒答",
"因果路径不足以支撑完整回答时\n直接说无法回答\n\n纯LLM从不承认自己不确定\nGraph RAG在55%的问题上\n选择诚实说明局限", ORANGE)
add_card(s, Inches(8.85), Inches(1.4), Inches(3.85), Inches(2.8), "每条陈述标注来源",
"格式: [路径1] [路径2,3]\n\n确保每个结论都能\n追溯到图谱中的具体因果路径\n\n让用户自己判断\n这个答案到底有多可靠", GREEN)
add_bottom(s, "约束机制本身是有效的 但约束之后的回答质量取决于图谱质量", Inches(4.8))

# ══════════════════ P10 双存储 ══════════════════
s = ns(sec=2)
add_title(s, "双存储: Neo4j + SQLite")
add_card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(4.7), "Neo4j 图数据库",
"因果链存储 + 变长路径查询\n6,976 节点  23,111 关系\nAccident聚合节点 1,579个\nUNIQUE约束: Equipment / Material / Accident\n\nA导致B这类关系\nNeo4j一行MERGE解决\nSQL里需要复杂自联结", ACCENT)
add_card(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(4.7), "SQLite 关系数据库",
"事故记录 1,579条 全部含根因后果\n化学品物性 72种 安全数据完整\n天气记录 108条\n4个索引加速查询\n\nCOUNT + GROUP BY 在SQL简洁\n在Cypher里反而绕\nDataLinker做跨源链接和一致性校验", GREEN)

# ══════════════════ P11 数据源 ══════════════════
s = ns(sec=2)
add_title(s, "三大数据源")
add_card(s, Inches(0.55), Inches(1.4), Inches(3.85), Inches(3.3), "mem.gov.cn",
"全量95个月度汇编页\n1,261份事故报告\n一事一报 平均150字\nBeautifulSoup HTML解析\n逐页提取事故标题和正文", ACCENT)
add_card(s, Inches(4.7), Inches(1.4), Inches(3.85), Inches(3.3), "微信公众号",
"化工安全教育服务平台\n108篇 清洗后保留74篇\n\n含安全建议和教训反思\nMitigation节点的主要来源\n月度简报没有应急措施段落\n这些信息来自微信文章", GREEN)
add_card(s, Inches(8.85), Inches(1.4), Inches(3.85), Inches(3.3), "PubChem + Open-Meteo",
"72种危化品物性\n闪点26种 爆炸极限31种 毒性43种\nCAS号 IUPAC名 分子量\n\nOpen-Meteo历史天气108条\n与事故地点日期匹配\n998条事故标注地理位置", ORANGE)
add_bottom(s, "注意: 主数据源mem简报每篇仅150字 这个特征对后续结论有决定性影响", Inches(5.3))

# ══════════════════ P12 分隔 ══════════════════
s = ns(dots=False)
add_section_page(s, "02", "实验验证", "对照实验 / 发现 / 根因分析")

# ══════════════════ P13 数据规模 ══════════════════
s = ns(sec=3)
add_title(s, "构建结果")
img = os.path.join(CHART_DIR, "01_node_types.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.55), Inches(1.3), Inches(6.2), Inches(5.2))
add_card(s, Inches(7.1), Inches(1.3), Inches(5.6), Inches(5.2), "数据统计",
"Mitigation: 4 → 71\n来源: 微信公众号事故分析文章\n\nAccident聚合节点: 0 → 1,579\nv0.7新功能 每起事故有独立聚合节点\n\nAbnormal占比: 67% → 51%\n事件原子化规则生效\n实体类型分布更均衡\n\n天气: 8 → 108\n化学品: 29 → 72\n\n图谱规模可观\n但节点越多就一定越好吗？", GREEN)

# ══════════════════ P14 实验设计 ══════════════════
s = ns(sec=3)
add_title(s, "对照实验: 图谱到底有没有用")
add_card(s, Inches(0.55), Inches(1.4), Inches(3.85), Inches(3.0), "关键词 RAG",
"jieba分词\nSQLite文本检索 Top 8\nLLM基于检索结果回答\n\n有来源 但检索噪声大\n无关文本可能干扰回答\n检索质量取决于关键词选择", ORANGE)
add_card(s, Inches(4.7), Inches(1.4), Inches(3.85), Inches(3.0), "Graph RAG",
"三层实体匹配\nCypher因果路径检索 max 4跳\n去重 + 子路径过滤\n约束生成 + [路径N]标注\n\n严格约束在图谱范围内\n但回答质量受图谱质量限制\n节点多时 无关因果链大量混入", ACCENT)
add_card(s, Inches(8.85), Inches(1.4), Inches(3.85), Inches(3.0), "纯 LLM",
"无外部数据\n仅凭模型参数知识回答\n\n没有任何约束\n引用的实体可能来自训练数据\n有些可能是对的\n但你无法区分来源\n\n从不拒答 每题引用5.95个图谱外实体", RED)
add_bottom(s, "20题 × 7种因果模式 × 3组baseline × 4维评估  同模型: DeepSeek v4-flash", Inches(4.9))

# ══════════════════ P15 实验结果 ══════════════════
s = ns(sec=3)
add_title(s, "实验结果")
img = os.path.join(CHART_DIR, "07_experiment.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.55), Inches(1.3), Inches(7.5), Inches(5.2))
add_card(s, Inches(8.4), Inches(1.3), Inches(4.35), Inches(5.2), "发现",
"约束效果 ✓\n图谱外实体: 5.95 → 0.65\n图内约束率: 5% → 70%\n诚实拒答率: 0% → 55%\n\n但实际使用效果不理想 ✗\n节点多导致检索慢\n无关因果链大量混入\n回答碎片化 看不清完整事故\n\n根本问题: 图谱质量不够\n不是约束机制的问题", RED)

# ══════════════════ P16 洞察1 ══════════════════
s = ns(sec=3)
add_title(s, "事故类型与时间趋势")
for img_file, left in [(os.path.join(CHART_DIR, "02_accident_types.png"), Inches(0.3)),
                        (os.path.join(CHART_DIR, "03_timeline.png"), Inches(6.7))]:
    if os.path.exists(img_file):
        s.shapes.add_picture(img_file, left, Inches(1.4), Inches(6.3), Inches(5.1))

# ══════════════════ P17 洞察2 ══════════════════
s = ns(sec=3)
add_title(s, "高频化学品与设备")
for img_file, left in [(os.path.join(CHART_DIR, "04_chemicals.png"), Inches(0.3)),
                        (os.path.join(CHART_DIR, "05_equipment.png"), Inches(6.7))]:
    if os.path.exists(img_file):
        s.shapes.add_picture(img_file, left, Inches(1.4), Inches(6.3), Inches(5.1))

# ══════════════════ P18 洞察3 ══════════════════
s = ns(sec=3)
add_title(s, "异常状态与14条统计发现")
img = os.path.join(CHART_DIR, "06_dangerous_states.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.3), Inches(1.4), Inches(6.3), Inches(5.1))
add_card(s, Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.1), "14条统计洞察",
"1  爆炸48%+中毒窒息22% = 事故的70%\n2  2010年代616起高峰 2020s降至128\n3  山东142 江苏88 河北72 Top3省份\n4  硫化氢49 氮气43 甲醇39 Top3\n5  反应釜62起 4倍于第二名的管道(16)\n6  违规操作31%+设备故障27% = 58%\n7  盲目施救166起(11%) 有限空间模式\n8  低闪点化学品事故频率2.8倍于高闪点\n9  电磁阀+氮气为最危险组合\n10 冬季爆炸占比(62%)略高于夏季(61%)\n11 4月事故最多(127) 9月最少(81)\n12 微信数据覆盖率高于mem简报\n13 标题爆炸出现651次(41%)\n14 2020s爆炸80起 中毒27起 同比大降\n\n这些发现来自统计分析\n不需要知识图谱也能得到", ACCENT)

# ══════════════════ P19 分隔 ══════════════════
s = ns(dots=False)
add_section_page(s, "03", "根因分析", "为什么图谱没有达到预期效果")

# ══════════════════ P20 根因 ══════════════════
s = ns(sec=3)
add_title(s, "为什么图谱没有达到预期")
add_card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(2.3), "数据源天花板",
"月度简报平均150字\n信息量本身就不足以支撑\n深层的因果推理\n\nLLM抽取的因果链最多3跳\n不是Prompt的问题\n是源材料没有提供更深的因果信息", RED)
add_card(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.3), "节点爆炸 信息杂糅",
"6,976个节点在检索时被大量激活\n无关因果链混入回答\n用户无法分辨哪些链条相关\n\n理想的因果图应该分叉汇聚\n但我们的图谱是碎片化的\n同义实体未合并加剧了这个问题", RED)
add_card(s, Inches(0.55), Inches(4.2), Inches(5.9), Inches(2.3), "响应速度退化",
"嵌入模型加载数秒\n6,000实体的相似度计算\n多层匹配的串行开销\nCypher变长路径遍历\n\n节点越多 响应越慢\n规模增长与可用性下降同步", ORANGE)
add_card(s, Inches(6.8), Inches(4.2), Inches(5.9), Inches(2.3), "统计分析反而更直接",
"事故类型 频次 排行 地域分布\n最有价值的洞察\nSQL统计就能得到\n不需要知识图谱\n\n知识图谱的优势在因果推理\n但在简报数据上\n这个优势发挥不出来", ORANGE)

# ══════════════════ P21 分隔 ══════════════════
s = ns(dots=False)
add_section_page(s, "04", "系统演示", "架构 / 问答 / 数据分析")

# ══════════════════ P22 演示 ══════════════════
s = ns(sec=4)
add_title(s, "系统演示")
add_card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(2.5), "问答演示",
"展示Graph RAG的约束效果\n\n演示问题: 硫化氢溢出会造成什么危害\n展示约束回答和来源引用\n\n再演示图谱信息不足时的拒答\n反应釜温度失控如何导致爆炸\n系统诚实说明无法回答", ACCENT)
add_card(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.5), "数据分析与可视化",
"展示6个分析标签页\n事故类型 趋势 化学品 设备分析\n14条统计洞察问答\n\n交互式图谱浏览\n展示节点规模和碎片化现象", GREEN)
add_bottom(s, "[现场演示]", Inches(4.6))

# ══════════════════ P23 分隔 ══════════════════
s = ns(dots=False)
add_section_page(s, "05", "收获与反思", "学到了什么 / 如果重来会怎么做")

# ══════════════════ P24 收获 ══════════════════
s = ns(sec=5)
add_title(s, "三个收获")
add_card(s, Inches(0.55), Inches(1.4), Inches(3.85), Inches(4.8), "工程能力",
"从零搭建了完整的五层系统\n数据采集 LLM抽取 图数据库\n检索匹配 前端可视化\n\nPrompt Chain迭代十几版\n从失败中归纳规则的工程方法\n比任何单一技术都更有价值\n\n双存储设计是有效的\n图数据库和关系数据库各司其职", ACCENT)
add_card(s, Inches(4.7), Inches(1.4), Inches(3.85), Inches(4.8), "实验方法",
"用对照实验来检验假设\n而不是凭感觉下结论\n\n三组baseline 同模型 同问题\n这个方法论本身是值得的\n即使结论推翻了初始假设\n\n实验设计比实验结果更重要\n它是可以迁移到其他项目的", GREEN)
add_card(s, Inches(8.85), Inches(1.4), Inches(3.85), Inches(4.8), "核心教训",
"知识图谱不是银弹\n\n数据质量决定图谱质量\n图谱质量决定回答质量\n\n对于短报式的事故数据\n知识图谱的投入产出比不合理\n约束有效 但被约束的内容不够好\n\n选择合适的数据源\n比优化算法更重要", RED)

# ══════════════════ P25 反思 ══════════════════
s = ns(sec=5)
add_title(s, "如果重来会怎么做")
add_card(s, Inches(0.55), Inches(1.4), Inches(5.9), Inches(2.3), "先确定数据质量再定方案",
"优先寻找更详细的数据源\n如CSB调查报告 每份数万字\n\n在数据源确定之前\n不急于构建大规模图谱\n\n先用小规模高质量数据\n验证核心假设\n再决定是否扩展", ACCENT)
add_card(s, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.3), "改进实验设计",
"对照实验的设计需要更严谨\n\n图谱外实体不等于错误\n纯LLM引用的实体\n有些可能确实是对的\n\n应该区分来源可追溯性\n和事实正确性\n而不是简单统计实体数量", GREEN)
add_card(s, Inches(0.55), Inches(4.2), Inches(5.9), Inches(2.3), "先做实体消歧 控制图谱规模",
"同义实体合并应该在抽取之后立即做\n形成爆炸性混合气体\n和形成爆炸性混合物应该自动合并\n\n按事故类型或化学品分类\n构建多个小型专用图谱\n小规模 高质量 可解释\n比大规模 低质量 碎片化更有用", ORANGE)
add_card(s, Inches(6.8), Inches(4.2), Inches(5.9), Inches(2.3), "建立自动化测试",
"十几版Prompt迭代\n全是手动跑几个例子看效果\n改了Prompt不知道是变好还是变坏\n\n人工标注20份以上黄金测试集\n每次改动自动回归\n优化效率提升5倍以上", PURPLE)

# ══════════════════ P26 分隔 ══════════════════
s = ns(dots=False)
add_section_page(s, "", "总结")

# ══════════════════ P27 总结 ══════════════════
s = ns(sec=0)
add_title(s, "总结")
add_card(s, Inches(0.55), Inches(1.4), Inches(2.85), Inches(4.5), "做了什么",
"构建了完整的KG+LLM问答系统\n1,579起事故 6,976节点 23,111边\n从数据采集到问答的全流程\n\n设计了三组对照实验\n验证了约束机制的有效性\n图谱外实体: 5.95 → 0.65\n诚实拒答率: 0% → 55%", ACCENT)
add_card(s, Inches(3.65), Inches(1.4), Inches(2.85), Inches(4.5), "发现了什么",
"约束有效 但回答质量不理想\n\n节点爆炸导致信息杂糅\n检索速度随规模增长而下降\n碎片化的因果链\n无法呈现真实的事故因果结构\n\n14条统计分析发现反而更有价值\n而这些不需要知识图谱", GREEN)
add_card(s, Inches(6.75), Inches(1.4), Inches(2.85), Inches(4.5), "为什么",
"源材料的质量决定了图谱的天花板\n150字简报无法支撑深层因果推理\n\n同义实体未合并加剧碎片化\n无关因果链在检索时大量混入\n\n知识图谱适合有详细因果描述的\n长文本数据\n不适合简报式的短文本数据", RED)
add_card(s, Inches(9.85), Inches(1.4), Inches(2.85), Inches(4.5), "如果重来",
"先确定数据源质量再决定技术方案\n\n用小规模高质量数据验证假设\n不急于构建大规模图谱\n\n改进实验设计\n区分来源可追溯性和事实正确性\n\n先做实体消歧\n按事故类型建专用小图", ORANGE)

# ══════════════════ P28 感谢 ══════════════════
s = ns(dots=False)
tb = s.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.3), Inches(2.5)).text_frame
tb.paragraphs[0].text = "感谢聆听"
tb.paragraphs[0].font.size = Pt(56); tb.paragraphs[0].font.bold = True
tb.paragraphs[0].font.color.rgb = ACCENT; tb.paragraphs[0].alignment = PP_ALIGN.CENTER
tb.add_paragraph().text = "欢迎提问"
tb.paragraphs[1].font.size = Pt(28); tb.paragraphs[1].font.color.rgb = DARK
tb.paragraphs[1].alignment = PP_ALIGN.CENTER

prs.save(OUT_PATH)
print(f"Done — {len(prs.slides)} slides")
